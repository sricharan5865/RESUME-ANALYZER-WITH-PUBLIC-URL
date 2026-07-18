import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Bright & Professional Light Color Palette
    bg_color = RGBColor(248, 250, 252)         # Very light slate gray
    card_color = RGBColor(255, 255, 255)       # Crisp white cards
    border_color = RGBColor(226, 232, 240)     # Subtle borders
    text_primary = RGBColor(15, 23, 42)        # Charcoal black
    text_secondary = RGBColor(71, 85, 105)     # Soft slate text
    
    accent_primary = RGBColor(79, 70, 229)     # Deep Indigo
    accent_info = RGBColor(14, 116, 144)       # Corporate Teal
    accent_success = RGBColor(22, 101, 52)     # Forest Green
    accent_warning = RGBColor(185, 28, 28)     # Crimson Red
    accent_alert = RGBColor(217, 119, 6)       # Amber Accent

    # Image Paths
    media_dir = r"C:\Users\sri charan\.gemini\antigravity\brain\41405d65-abc2-422a-a370-27ab1394f687"
    img1_path = os.path.join(media_dir, "media__1783665147653.png")
    img2_path = os.path.join(media_dir, "media__1783665147668.png")
    img3_path = os.path.join(media_dir, "media__1783665147683.png")

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_card(slide, left, top, width, height, fill_color, line_color=border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
        return shape

    def add_header(slide, title, category="TALENTFLOW CAPABILITIES"):
        # Category label
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = accent_primary
        p.font.name = 'Arial'

        # Main Slide Title
        tx_box_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
        tf_title = tx_box_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = text_primary
        p_title.font.name = 'Arial'

    slide_layout = prs.slide_layouts[6]

    # ==========================================
    # --- SLIDE 1: Title Slide (Full Bleed Card) ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_card(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5), card_color, accent_primary)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TALENTFLOW"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p.font.name = 'Arial'
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Recruitment Automation Platform"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Streamlining the entire hiring lifecycle — from resume sourcing to candidate evaluation"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(16)
    p3.font.color.rgb = text_secondary
    p3.space_before = Pt(12)

    p4 = tf.add_paragraph()
    p4.text = "Generative AI  •  Semantic Search  •  Email Automation  •  Smart Match & Interview Generation"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = accent_info
    p4.space_before = Pt(24)

    # ==========================================
    # --- SLIDE 2: Problem Statement ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "The Problem Statement: Recruitment Bottlenecks", "THE RECRUITING CHALLENGE")
    
    # 2 Stretched Content Cards (Left & Right)
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_warning)
    tb_prob1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_prob1 = tb_prob1.text_frame
    tf_prob1.word_wrap = True
    p = tf_prob1.paragraphs[0]
    p.text = "OPERATIONAL INEFFICIENCY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_warning
    
    p2 = tf_prob1.add_paragraph()
    p2.text = "Manual Screening & Fragmentation"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Time-Intensive Screening: Recruiters spend an average of 23 hours reviewing resumes per hire, creating major delays.",
        "Inbox Overload: Resumes arrive scattered across multiple email servers (Gmail, Outlook). Downloading and sorting attachments manually is slow.",
        "Poor Pipeline Tracking: Managing candidate progress in spreadsheets leads to missed dates, lost CVs, and slow communication."
    ]
    for b in bullets:
        pb = tf_prob1.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_warning)
    tb_prob2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_prob2 = tb_prob2.text_frame
    tf_prob2.word_wrap = True
    p = tf_prob2.paragraphs[0]
    p.text = "TECHNOLOGICAL LIMITS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_warning
    
    p2 = tf_prob2.add_paragraph()
    p2.text = "Keyword Matchers & Biased Prep"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Keyword-Only Search: Traditional ATS systems miss candidates due to exact word queries (e.g. 'Node.js Developer' ignoring 'Express.js Engineer').",
        "Inconsistent Evaluation: No objective scoring exists. Candidates are interviewed on subjective standards without structured checklists.",
        "Unprepared Interview Panels: Hiring managers ask mismatched or generic questions, failing to test candidate skills effectively."
    ]
    for b in bullets:
        pb = tf_prob2.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # ==========================================
    # --- SLIDE 3: The Solution ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "The Solution: Resilient Recruitment Intelligence", "TALENTFLOW ADVANTAGE")
    
    # 2 Stretched Content Cards (Left & Right)
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_success)
    tb_sol1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_sol1 = tb_sol1.text_frame
    tf_sol1.word_wrap = True
    p = tf_sol1.paragraphs[0]
    p.text = "AUTOMATED CHANNELS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_success
    
    p2 = tf_sol1.add_paragraph()
    p2.text = "Sourcing, Parsing & OCR"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Dynamic Email Pollers: Continual background sync scans Gmail (IMAP) & Outlook (Graph API) to capture resumes automatically.",
        "AI Parsing & OCR Ingest: Parses raw PDF/Docx text into database structures, falling back to Python Tesseract OCR for scanned files.",
        "Duplicate Resolution Flow: Prevents candidate database clutter via a 4-option duplicate candidate management flow."
    ]
    for b in bullets:
        pb = tf_sol1.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_success)
    tb_sol2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_sol2 = tb_sol2.text_frame
    tf_sol2.word_wrap = True
    p = tf_sol2.paragraphs[0]
    p.text = "INTELLIGENT RECRUITING"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_success
    
    p2 = tf_sol2.add_paragraph()
    p2.text = "RAG, JD Score & Interview Prep"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Semantic Search (RAG): Maps queries to vectors for conceptual match, allowing searches like 'Node.js' to find 'Express.js' profiles.",
        "Objective JD Match: Ranks candidates (0-100%) against JDs with detailed missing/matching skill metrics and AI reasoning.",
        "Custom Question Banks: Auto-generates HR & technical question guides with correct answers, sorted by criticality (Must Ask, optional)."
    ]
    for b in bullets:
        pb = tf_sol2.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # ==========================================
    # --- SLIDE 4: Sourcing Channels (Input Methods) ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Candidate Sourcing: Automated Input Channels", "INPUT METHODS")
    
    sourcing_channels = [
        ("Direct Upload", "Allows recruiters to drag and drop single or bulk resume files. Accepts PDF & DOCX formats. Immediately runs text parsing and database indexing.", accent_primary),
        ("Gmail IMAP Sync", "Background sync loop polls Gmail folders every 30s using app password connection. Downloads PDF attachments and logs processed UIDs to avoid double imports.", accent_warning),
        ("Outlook API Connection", "Secured Microsoft Graph OAuth 2.0 flow. Fetches unread emails, parses attachments, and enables direct platform email replies.", accent_info),
        ("Auto Poller Loop", "A unified background poller checks settings to run the correct email pipeline, categorizing spam out and auto-importing applicants.", accent_success)
    ]
    for i, (title, desc, color) in enumerate(sourcing_channels):
        col_left = Inches(0.5) + i * (Inches(3.0) + Inches(0.111))
        add_card(slide, col_left, Inches(1.4), Inches(3.0), Inches(5.6), card_color, color)
        tb = slide.shapes.add_textbox(col_left + Inches(0.15), Inches(1.6), Inches(2.7), Inches(5.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"METHOD 0{i+1}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(19)
        p2.font.bold = True
        p2.font.color.rgb = text_primary
        p2.space_before = Pt(10)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = text_secondary
        p3.space_before = Pt(14)

    # ==========================================
    # --- SLIDE 5: Duplicate Resolution Flow ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Ingestion Control & Duplicate Candidate Flow", "DATA CONFLICT RESOLUTION")
    
    # Left Card - Conflict Detection
    add_card(slide, Inches(0.5), Inches(1.4), Inches(5.0), Inches(5.6), card_color)
    tb_det = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4.6), Inches(5.2))
    tf_det = tb_det.text_frame
    tf_det.word_wrap = True
    p = tf_det.paragraphs[0]
    p.text = "INTELLIGENT DEDUPLICATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_alert
    
    p2 = tf_det.add_paragraph()
    p2.text = "How Conflicts are Spotted"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Email Verification: Matches incoming candidates against database records using case-insensitive email regex checking.",
        "Name Matching: Employs case-insensitive name matching to catch candidates applying with alternative email accounts.",
        "Atomic Upserts: Prevents simultaneous ingestion write conflicts via MongoDB atomic checks."
    ]
    for b in bullets:
        pb = tf_det.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(15)

    # Right Card - Four Resolution Options
    add_card(slide, Inches(5.8), Inches(1.4), Inches(7.0), Inches(5.6), card_color, accent_primary)
    tb_opt = slide.shapes.add_textbox(Inches(6.0), Inches(1.6), Inches(6.6), Inches(5.2))
    tf_opt = tb_opt.text_frame
    tf_opt.word_wrap = True
    
    p = tf_opt.paragraphs[0]
    p.text = "CONFLICT DECISION FLOW"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    
    p2 = tf_opt.add_paragraph()
    p2.text = "Four Resolution Options Offered to User"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    options = [
        ("1. Update (Overwrite Existing Info & CV)", "Retains current Candidate ID, updates experience, education, skills, and replaces PDF resume file with new upload. Triggers scoring re-runs."),
        ("2. Delete Existing & Import New", "Wipes the previous profile (including history and indices) and imports the new resume fresh, assigning a new Candidate ID."),
        ("3. Delete Existing Only (Halt Import)", "Deletes the existing candidate from the database entirely, while discarding the newly uploaded file without importing it."),
        ("4. Cancel (Discard Uploaded File)", "Keeps the existing candidate database profile untouched, immediately deleting the new temporary file from the server.")
    ]
    for opt_title, opt_desc in options:
        pb_title = tf_opt.add_paragraph()
        pb_title.text = opt_title
        pb_title.font.size = Pt(14)
        pb_title.font.bold = True
        pb_title.font.color.rgb = accent_info
        pb_title.space_before = Pt(10)
        
        pb_desc = tf_opt.add_paragraph()
        pb_desc.text = opt_desc
        pb_desc.font.size = Pt(12)
        pb_desc.font.color.rgb = text_secondary
        pb_desc.space_before = Pt(2)

    # ==========================================
    # --- SLIDE 6: Smart Resume Ingestion & Parsing ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Smart Resume Ingestion: AI & OCR Parser", "PROFILE INGESTION")
    
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color)
    tb_prs = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_prs = tb_prs.text_frame
    tf_prs.word_wrap = True
    p = tf_prs.paragraphs[0]
    p.text = "DOCUMENT EXTRACTION PIPELINE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    
    p2 = tf_prs.add_paragraph()
    p2.text = "Resume Parsing Mechanics"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Multi-Format Parser: Extracts textual details from PDF, DOCX, TXT, RTF, and MD files.",
        "PDF.js Core: Extracts text segments directly using web-standard client libraries for high speed.",
        "Tesseract OCR Fallback: Runs a Python subprocess utilizing Tesseract to extract content if the PDF is scanned or image-based.",
        "JSON Repair Engine: Employs regex filters and structural bracket matching to repair malformed LLM outputs and avoid crashes."
    ]
    for b in bullets:
        pb = tf_prs.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # Code card
    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, border_color)
    tb_code = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.6), Inches(5.0))
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    p = tf_code.paragraphs[0]
    p.text = "Extracted Candidate Schema Model:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p2 = tf_code.add_paragraph()
    p2.text = "{\n  \"name\": \"John Doe\",\n  \"email\": \"john@example.com\",\n  \"skills\": [\"React\", \"Node.js\", \"AWS\"],\n  \"experience\": [\n    {\n      \"role\": \"Backend Engineer\",\n      \"company\": \"TechCorp\",\n      \"duration\": \"3 years 4 months\"\n    }\n  ],\n  \"seniorityLevel\": \"Senior\",\n  \"tags\": [{\"value\": \"Backend\", \"category\": \"Domain\"}]\n}"
    p2.font.size = Pt(12)
    p2.font.name = "Consolas"
    p2.font.color.rgb = accent_info
    p2.space_before = Pt(15)

    # ==========================================
    # --- SLIDE 7: Multi-Provider AI Engine ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Multi-Provider AI & Model Settings", "AI SERVICE CONFIGURATION")
    
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color)
    tb_ai = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_ai = tb_ai.text_frame
    tf_ai.word_wrap = True
    p = tf_ai.paragraphs[0]
    p.text = "PLUG-AND-PLAY AI LAYERS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    
    p2 = tf_ai.add_paragraph()
    p2.text = "Supported Models & APIs"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Google Gemini: Connects directly to gemini-2.0-flash using standard developer API keys.",
        "OpenAI: Dispatches to GPT-4o for parsing and evaluations.",
        "Anthropic Claude: Integrated with Claude 3.5 Sonnet.",
        "OpenRouter: Allows routing requests to 100+ cloud LLMs via a single sk-or-* prefixed API key.",
        "Ollama (Local LLM): Connects to local endpoints (e.g. LLaMA3) for offline usage and data confidentiality."
    ]
    for b in bullets:
        pb = tf_ai.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # Connection Checks Card
    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_info)
    tb_conn = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_conn = tb_conn.text_frame
    tf_conn.word_wrap = True
    
    p = tf_conn.paragraphs[0]
    p.text = "ROBUSTNESS & OPTIMIZATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_info
    
    p2 = tf_conn.add_paragraph()
    p2.text = "AI Operations & Stability"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets_op = [
        "Connection Checks: Settings UI offers test-connection endpoints for Ollama, Gmail, and Outlook to verify integrations before saving.",
        "Ollama Context Optimization: Dynamically shrinks local prompt schemas to avoid high pre-processing latency.",
        "High Output Limits: Enforces at least 8000 max_tokens for candidate analysis & question generators to prevent JSON truncation.",
        "Stateless Fallbacks: If AI APIs fail, fallback routines preserve database status without losing uploaded resume files."
    ]
    for b in bullets_op:
        pb = tf_conn.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(13.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # ==========================================
    # --- SLIDE 8: Advanced Semantic Search (RAG) ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Advanced Semantic Search (RAG Engine)", "AI SEARCH SYSTEM")
    
    # Left Card
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color)
    tb_rag = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_rag = tb_rag.text_frame
    tf_rag.word_wrap = True
    p = tf_rag.paragraphs[0]
    p.text = "CONCEPTUAL SEARCH RETRIEVAL"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = accent_info
    
    p2 = tf_rag.add_paragraph()
    p2.text = "RAG & Vector Embeddings"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Chunk-Level Indexing: Splits resumes into sections (contact, skills, experience, projects) to store embeddings in MongoDB.",
        "Semantic Vector Search: Translates natural language queries into embeddings to locate contextually relevant candidate chunks.",
        "AI Q&A Chatbot: An interactive 'Ask AI' chat panel answers recruiter questions using relevant candidates as context.",
        "Keyword Fallback: Seamless token matching checks for exact keyword matches if the vector API goes offline."
    ]
    for b in bullets:
        pb = tf_rag.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # Right Embedding Graph/Mockup
    if os.path.exists(img2_path):
        slide.shapes.add_picture(img2_path, Inches(6.8), Inches(1.4), width=Inches(6.0), height=Inches(5.6))
    else:
        add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, border_color)
        tb_img = slide.shapes.add_textbox(Inches(6.8), Inches(3.5), Inches(6.0), Inches(1.0))
        tb_img.text_frame.text = "[RAG Search Visual / Graph]"
        tb_img.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb_img.text_frame.paragraphs[0].font.color.rgb = text_secondary

    # ==========================================
    # --- SLIDE 9: High-Speed Tag & Keyword Search ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "High-Speed Tag & Keyword Search", "SEARCH METHODOLOGIES")
    
    # 3 Stretched cards (filling spaces)
    search_types = [
        ("Trie Autocomplete", "Indexes candidate tags and skills for instantaneous keystroke suggestions.", accent_primary),
        ("Trigram Fuzzy Match", "Splits words into three-character tokens to match queries despite spelling mistakes.", accent_info),
        ("Interactive Tag Cloud", "Visualized tag list for rapid exploration of pool skills and counts.", accent_success)
    ]
    for i, (title, desc, color) in enumerate(search_types):
        col_left = Inches(0.5) + i * (Inches(4.0) + Inches(0.166))
        add_card(slide, col_left, Inches(1.4), Inches(4.0), Inches(5.6), card_color)
        tb = slide.shapes.add_textbox(col_left + Inches(0.2), Inches(1.6), Inches(3.6), Inches(5.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"METHOD 0{i+1}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = text_primary
        p2.space_before = Pt(15)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(14)
        p3.font.color.rgb = text_secondary
        p3.space_before = Pt(12)

    # ==========================================
    # --- SLIDE 10: Job Description (JD) Matcher ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Job Description (JD) Matcher & Reranking", "CANDIDATE RANKING")
    
    # Stretched left card
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color)
    tb_jd = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_jd = tb_jd.text_frame
    tf_jd.word_wrap = True
    p = tf_jd.paragraphs[0]
    p.text = "INTELLIGENT RANKING ALGORITHM"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = accent_success
    p2 = tf_jd.add_paragraph()
    p2.text = "JD Scoring & Analysis"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Upload Job Description: Enter title, location, requirements, and description.",
        "Scoring Scale (0-100%): Ranks profiles objectively using AI models.",
        "Dual Scoring Modes: JD Match ranks candidates for a selected job, while Own-Category Score ranks them within their parsed specialization.",
        "Retrieve-and-Rerank: Uses vector search to pre-filter top candidates, scoring only matches via LLM to optimize speed and API costs."
    ]
    for b in bullets:
        pb = tf_jd.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(13.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(12)

    # Stretched right card
    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, border_color)
    tb_vis = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_vis = tb_vis.text_frame
    tf_vis.word_wrap = True
    p = tf_vis.paragraphs[0]
    p.text = "JD Matching Logic Flow:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p2 = tf_vis.add_paragraph()
    p2.text = "1. Parse Job Description requirements\n\n2. Query RAG vector database for top candidates\n\n3. Perform programmatical experience duration validation\n\n4. Dispatch selected profiles to LLM for precise match score\n\n5. Extract matching skills, missing gaps, and reasoning"
    p2.font.size = Pt(14)
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(20)

    # ==========================================
    # --- SLIDE 11: Interactive Kanban Pipeline Board ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Interactive Kanban Pipeline Board", "PIPELINE MANAGEMENT")
    
    # 5 Stretched pipeline columns (minimal side space) with detailed feature bullets
    stages = [
        (
            "Inbox", 
            "New email & upload profiles.", 
            accent_primary,
            ["• Sourced email CVs", "• Manual file uploads", "• Duplicate prevention", "• OCR text fallback"]
        ),
        (
            "Shortlist", 
            "Selected fit profiles.", 
            accent_info,
            ["• Experience math validation", "• Skill checks", "• Initial candidate review", "• Custom tags display"]
        ),
        (
            "Interview", 
            "Scheduled panels & tests.", 
            accent_warning,
            ["• Custom technical Q&As", "• Custom HR Q&As", "• Regenerable question bank", "• Checklists for interviewers"]
        ),
        (
            "Offered", 
            "Fit profiles for hiring.", 
            accent_success,
            ["• Profile score matching", "• Fit recommendations", "• Final review pipeline", "• Offer letter template"]
        ),
        (
            "Rejected", 
            "Archived applicant files.", 
            text_secondary,
            ["• Archive with reason", "• Searchable pool retention", "• Re-engagement rules", "• Timeline logs"]
        )
    ]
    for i, (stage_name, desc, color, bullets) in enumerate(stages):
        col_left = Inches(0.5) + i * (Inches(2.4) + Inches(0.066))
        add_card(slide, col_left, Inches(1.4), Inches(2.4), Inches(5.6), card_color, color)
        tb = slide.shapes.add_textbox(col_left + Inches(0.1), Inches(1.5), Inches(2.2), Inches(5.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = stage_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = text_primary
        p2.space_before = Pt(8)
        
        for bullet in bullets:
            pb = tf.add_paragraph()
            pb.text = bullet
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = text_secondary
            pb.space_before = Pt(8)
            
        p3 = tf.add_paragraph()
        p3.text = "History & Audit Logs"
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = color
        p3.space_before = Pt(16)

    # ==========================================
    # --- SLIDE 12: Details Drawer & PDF Preview ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Slide-Out Details Drawer & PDF Preview", "RECRUITER WORKSPACE")
    
    # Left context workspace
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color)
    tb_drw = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_drw = tb_drw.text_frame
    tf_drw.word_wrap = True
    p = tf_drw.paragraphs[0]
    p.text = "CONTEXT RETAINED INTERFACE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p2 = tf_drw.add_paragraph()
    p2.text = "360° Candidate Profile View"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Employment Gap Detection: Automatically calculates total experience and logs detailed employment gap timelines.",
        "Interactive Side-by-Side: View candidate's original PDF/Docx alongside AI-structured data.",
        "Tabbed Detail Panels: Organizes structured data into Experience, Education, Projects, and customized Question-Answers.",
        "Audit Log Trail: Lists detailed profile history entries (Ingested, Stage Changed, Emailed)."
    ]
    for b in bullets:
        p_b = tf_drw.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = text_secondary
        p_b.space_before = Pt(14)

    # Right slide image
    if os.path.exists(img3_path):
        slide.shapes.add_picture(img3_path, Inches(6.8), Inches(1.4), width=Inches(6.0), height=Inches(5.6))
    else:
        add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, border_color)
        tb_img = slide.shapes.add_textbox(Inches(6.8), Inches(3.5), Inches(6.0), Inches(1.0))
        tb_img.text_frame.text = "[Slide-Out Details Drawer Mockup]"
        tb_img.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb_img.text_frame.paragraphs[0].font.color.rgb = text_secondary

    # ==========================================
    # --- SLIDE 13: Tailored Interview Question Generator ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Tailored Interview Question Generator", "PRE-SCREENING PREPARATION")
    
    # 2 Content Cards
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_primary)
    tb_q1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_q1 = tb_q1.text_frame
    tf_q1.word_wrap = True
    p = tf_q1.paragraphs[0]
    p.text = "CANDIDATE-SPECIFIC CHECKS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p2 = tf_q1.add_paragraph()
    p2.text = "HR & Technical Q&A Generation"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Culture & Behavioral Questions: Generates custom HR screening questions aligned with candidate seniority.",
        "Domain-Specific Questions: Inspects candidate skills to build custom technical evaluation questions.",
        "Expected Answer Guide: Supplies the correct expected answers for each question so any interviewer can grade responses.",
        "Regenerate on Demand: A simple button in the details panel lets users re-run AI generation to fresh sets."
    ]
    for b in bullets:
        p3 = tf_q1.add_paragraph()
        p3.text = "• " + b
        p3.font.size = Pt(14)
        p3.font.color.rgb = text_secondary
        p3.space_before = Pt(14)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_info)
    tb_q2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_q2 = tb_q2.text_frame
    tf_q2.word_wrap = True
    p = tf_q2.paragraphs[0]
    p.text = "COMPETENCY HIGHLIGHTING"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_info
    p2 = tf_q2.add_paragraph()
    p2.text = "Importance Ratings & Badging"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    ratings = [
        ("🔴 Must Ask / Very Important", "Critical questions targeting core skill gaps or mismatch points."),
        ("🟠 Important", "Highlights crucial job description requirements."),
        ("🟡 Good to Ask", "Helpful to confirm details, but not project blockers."),
        ("🔵 Screening", "Rapid qualification checks (e.g. visa status, notice period)."),
        ("🟢 Optional", "Nice-to-have topics if interview time permits.")
    ]
    for r_title, r_desc in ratings:
        pb_title = tf_q2.add_paragraph()
        pb_title.text = r_title
        pb_title.font.size = Pt(13)
        pb_title.font.bold = True
        pb_title.font.color.rgb = text_primary
        pb_title.space_before = Pt(8)
        
        pb_desc = tf_q2.add_paragraph()
        pb_desc.text = r_desc
        pb_desc.font.size = Pt(11.5)
        pb_desc.font.color.rgb = text_secondary
        pb_desc.space_before = Pt(1)

    # ==========================================
    # --- SLIDE 14: Operations: Inbox & Ingestion Log ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Inbox Management & Real-Time Ingestion Logs", "RECRUITER WORKFLOW")
    
    # Left Card
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_primary)
    tb_inbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_inbox = tb_inbox.text_frame
    tf_inbox.word_wrap = True
    p = tf_inbox.paragraphs[0]
    p.text = "COMMUNICATION HUB"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p2 = tf_inbox.add_paragraph()
    p2.text = "Unified Sourcing Inbox"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Email Previews: Displays sender, subject, date, body snippet, and attachments directly in the app.",
        "AI Email Categorization: Labels emails into Resume, HR, Spam, Client, Interview, Notification, or Other.",
        "Email Send Modal: Compose and send messages to candidates with templates (Interview, Offer, Rejection) using mail servers.",
        "Auto-Process Toggle: Automatically ingests emails categorized as 'Resume' without human intervention."
    ]
    for b in bullets:
        pb = tf_inbox.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(13.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # Right Card
    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_info)
    tb_track = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_track = tb_track.text_frame
    tf_track.word_wrap = True
    p = tf_track.paragraphs[0]
    p.text = "INGESTION TRACKER"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_info
    p2 = tf_track.add_paragraph()
    p2.text = "Real-Time Processing Logs"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "Status Monitoring: Real-time poller tracks uploads through states: Processing, Success, Failed, Duplicate, Cancelled.",
        "Error Diagnostics: Log details display the exact failure reasons (e.g. failed OCR, missing API keys) for debugging.",
        "Candidate Link: Clickable success links instantly open the parsed candidate details drawer.",
        "Channel Source: Traces manual upload vs Gmail vs Outlook sourcing pathways."
    ]
    for b in bullets:
        pb = tf_track.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(13.5)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # ==========================================
    # --- SLIDE 15: Export Features & Data Management ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Data Export Features & Formats", "DATA MANAGEMENT")
    
    # 3 Stretched cards
    export_cards = [
        ("Pipeline CSV Export", "Allows exporting filtered candidates based on stages and date ranges. Generates UTF-8 BOM CSV files with parsed work experience, education, scores, and skills.", accent_primary),
        ("Tag Search CSV Export", "Directly exports candidate search matches filtered by tag queries and seniority sliders to spreadsheets.", accent_info),
        ("Ingestion Log CSV", "Exports tracker audits including filename, timestamp, sourcing channel, status, and failure logs for compliance audits.", accent_success)
    ]
    for i, (title, desc, color) in enumerate(export_cards):
        col_left = Inches(0.5) + i * (Inches(4.0) + Inches(0.166))
        add_card(slide, col_left, Inches(1.4), Inches(4.0), Inches(5.6), card_color, color)
        tb = slide.shapes.add_textbox(col_left + Inches(0.2), Inches(1.6), Inches(3.6), Inches(5.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"FORMAT 0{i+1}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = text_primary
        p2.space_before = Pt(15)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(14)
        p3.font.color.rgb = text_secondary
        p3.space_before = Pt(12)

    # ==========================================
    # --- SLIDE 16: Authentication & Security (RBAC) ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Authentication & Role-Based Access Control", "SECURITY & RBAC")
    
    # Left Card
    add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_primary)
    tb_sec = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_sec = tb_sec.text_frame
    tf_sec.word_wrap = True
    p = tf_sec.paragraphs[0]
    p.text = "JWT AUTHENTICATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p2 = tf_sec.add_paragraph()
    p2.text = "Session & Cryptography"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    bullets = [
        "JSON Web Tokens (JWT): Secures communications via token verification in headers or query params, set to expire in 24 hours.",
        "BCrypt Hashing: Encrypts passwords before storage in MongoDB to protect accounts.",
        "Candidate Sharing: Recruiters and admins can share specific candidate profiles with hiring managers via email assignment."
    ]
    for b in bullets:
        pb = tf_sec.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(14)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(14)

    # Right Card
    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), card_color, accent_info)
    tb_rbac = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_rbac = tb_rbac.text_frame
    tf_rbac.word_wrap = True
    p = tf_rbac.paragraphs[0]
    p.text = "RBAC PERMISSIONS MATRIX"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_info
    p2 = tf_rbac.add_paragraph()
    p2.text = "User Roles & Permissions"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(5)
    
    roles = [
        ("1. Administrator (Admin)", "Complete workspace authority. Configures system settings, toggles email syncs, adds users, and manages databases."),
        ("2. Recruiter", "Functional sourcing role. Can import profiles, trigger JD matches, update pipelines, compose emails, and run searches."),
        ("3. Hiring Manager", "Restricted reviewer. Access is limited to candidate profiles explicitly assigned to them. Cannot edit, upload, or delete data.")
    ]
    for r_title, r_desc in roles:
        pb_title = tf_rbac.add_paragraph()
        pb_title.text = r_title
        pb_title.font.size = Pt(13.5)
        pb_title.font.bold = True
        pb_title.font.color.rgb = accent_primary
        pb_title.space_before = Pt(10)
        
        pb_desc = tf_rbac.add_paragraph()
        pb_desc.text = r_desc
        pb_desc.font.size = Pt(12)
        pb_desc.font.color.rgb = text_secondary
        pb_desc.space_before = Pt(2)

    # ==========================================
    # --- SLIDE 17: Reporting & Analytics ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Reporting & Recruitment Analytics", "METRICS & KPI MONITOR")
    
    # 4 metrics cards
    kpis = [
        ("Conversion Funnel", "Visualizes conversion ratios from Sourced → Shortlisted → Interviewed → Offered.", accent_primary, Inches(0.5), Inches(1.4)),
        ("Match Score Distribution", "Segments the database into Excellent (90-100), Good (70-89), Average, and Poor bands.", accent_info, Inches(6.8), Inches(1.4)),
        ("Sourcing Channels", "Analyzes the intake volume of manual uploads vs Gmail vs Outlook integrations.", accent_success, Inches(0.5), Inches(4.2)),
        ("Talent Seniority Segments", "Graphs candidate pools across Junior, Mid, Senior, Lead, and Executive levels.", accent_warning, Inches(6.8), Inches(4.2))
    ]
    for title, desc, color, left, top in kpis:
        add_card(slide, left, top, Inches(6.0), Inches(2.6), card_color, color)
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.6), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13.5)
        p2.font.color.rgb = text_secondary
        p2.space_before = Pt(8)

    # ==========================================
    # --- SLIDE 18: Platform Summary ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_header(slide, "Platform Summary", "THE END-TO-END SUITE")
    
    # Stretched 4-card matrix
    summary_cards = [
        ("Automated Sourcing", "Monitors mailboxes dynamically for applicant resume attachments.", accent_primary, Inches(0.5), Inches(1.4)),
        ("Structured Parsing", "Hybrid PDF parser & OCR extractor structure documents cleanly.", accent_info, Inches(6.8), Inches(1.4)),
        ("Intelligent Search", "Semantic RAG & high-speed keyword systems optimize discovery.", accent_success, Inches(0.5), Inches(4.2)),
        ("Pipeline Management", "Job Description matching, custom Q&A arrays, and visual Kanban flow.", accent_warning, Inches(6.8), Inches(4.2))
    ]
    for title, desc, color, left, top in summary_cards:
        add_card(slide, left, top, Inches(6.0), Inches(2.6), card_color, color)
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.6), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = text_secondary
        p2.space_before = Pt(10)

    # ==========================================
    # --- SLIDE 19: Thank You & Call to Action (Full Bleed Slide) ---
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_card(slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5), card_color, accent_primary)
    
    concl_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf = concl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = accent_primary
    p.font.name = 'Arial'
    
    p2 = tf.add_paragraph()
    p2.text = "Transform Recruitment with Intelligent Automation"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = text_primary
    p2.space_before = Pt(15)

    p3 = tf.add_paragraph()
    p3.text = "Simplify Sourcing, Evaluation, and Pipeline Management.\nExplore the platform capabilities today."
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(15)
    p3.font.color.rgb = text_secondary
    p3.space_before = Pt(20)

    output_filename = "TalentFlow_Core_Features_Presentation.pptx"
    prs.save(output_filename)
    print(f"Success: {output_filename} has been created successfully!")

if __name__ == "__main__":
    create_presentation()
